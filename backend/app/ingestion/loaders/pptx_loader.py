from pathlib import Path

from pptx import Presentation

from app.ingestion.loaders.office_common import enforce_text_limit, validate_office_package
from app.models.document import Document, DocumentMetadata, FileType, ParsedUnit, SourceLocator, create_stable_document_id


class PPTXLoader:
    def load(self, file_path: str | Path) -> tuple[Document, list[ParsedUnit]]:
        path = Path(file_path)
        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Unsupported file type: {path.suffix}. Expected .pptx")
        content = validate_office_package(path, "ppt/presentation.xml")
        try:
            presentation = Presentation(path)
        except Exception as exc:
            raise ValueError(f"Unable to read PPTX: {path.name}") from exc
        units = []
        extracted = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title else None
            blocks = []
            for shape in slide.shapes:
                if shape.has_table:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    rows = [row for row in rows if any(row)]
                    if rows:
                        headers = [value or f"Column {index + 1}" for index, value in enumerate(rows[0])]
                        blocks.append("Columns: " + " | ".join(headers))
                        for row in rows[1:]:
                            blocks.append("Row:\n" + "\n".join(
                                f"{header}: {value}" for header, value in zip(headers, row) if value
                            ))
                elif getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text not in blocks:
                            blocks.append(("- " if paragraph.level else "") + text)
            rendered = "\n\n".join(blocks)
            if not rendered:
                continue
            extracted = enforce_text_limit(extracted, rendered)
            units.append(ParsedUnit(
                text=rendered,
                page_number=slide_number,
                heading=title,
                source_locator=SourceLocator(type="slide", label=f"Slide {slide_number}"),
            ))
        if not units:
            raise ValueError("PPTX contains no extractable text")
        properties = presentation.core_properties
        return Document(
            id=create_stable_document_id(content, FileType.PPTX),
            filename=path.name,
            file_type=FileType.PPTX,
            metadata=DocumentMetadata(
                title=properties.title or None,
                author=properties.author or None,
                page_count=len(presentation.slides),
            ),
        ), units
