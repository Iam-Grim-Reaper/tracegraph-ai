from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest
import re
from zipfile import ZIP_DEFLATED, ZipFile

from app.ingestion.service import IngestionService
from app.ingestion.loaders.xlsx_loader import XLSXLoader
from app.models.document import FileType


def create_docx(path):
    document = DocxDocument()
    document.add_heading("Enterprise Architecture", level=1)
    document.add_paragraph("Aster Analytics uses Apache Spark for batch data processing.")
    document.add_paragraph("The Data Platform team owns the Orion pipeline.")
    table = document.add_table(rows=2, cols=3)
    for cell, value in zip(table.rows[0].cells, ("System", "Owner", "SLA")):
        cell.text = value
    for cell, value in zip(table.rows[1].cells, ("Orion", "Data Platform", "4 hours")):
        cell.text = value
    document.save(path)


def create_pptx(path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Project Atlas"
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Technology"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(2))
    box.text_frame.text = "Atlas uses retrieval-augmented generation."
    box.text_frame.add_paragraph().text = "The Knowledge Systems team maintains Atlas."
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(1)).table
    for cell, value in zip(table.rows[0].cells, ("System", "Owner")):
        cell.text = value
    for cell, value in zip(table.rows[1].cells, ("Atlas", "Knowledge Systems")):
        cell.text = value
    presentation.save(path)


def create_xlsx(path, formula=False):
    workbook = Workbook()
    employees = workbook.active
    employees.title = "Employees"
    employees.append(["Name", "Role", "Skills", "Team"])
    employees.append(["Alice Morgan", "Data Engineer", "Python, Spark", "Data Platform"])
    employees.append([])
    employees.append(["Jordan Lee", "ML Engineer", "PyTorch, RAG", "AI Systems"])
    systems = workbook.create_sheet("Systems")
    systems.append(["System", "Owner", "Dependency"])
    systems.append(["Atlas", "AI Systems", "=1+1" if formula else "Qdrant"])
    systems.append(["Orion", "Data Platform", "Kafka"])
    workbook.save(path)


def create_dimensionless_xlsx(path):
    workbook = Workbook()
    employees = workbook.active
    employees.title = "Employees"
    employees.append(["Name", "Start date"])
    employees.append(["Aster", "2026-01-01"])
    workbook.save(path)

    with ZipFile(path) as archive:
        entries = {
            entry.filename: archive.read(entry.filename)
            for entry in archive.infolist()
        }
    entries["xl/worksheets/sheet1.xml"] = re.sub(
        br"<dimension[^>]*/>",
        b"",
        entries["xl/worksheets/sheet1.xml"],
    )
    with ZipFile(path, "w", ZIP_DEFLATED) as archive:
        for filename, content in entries.items():
            archive.writestr(filename, content)


@pytest.mark.parametrize(
    ("filename", "creator", "file_type"),
    [
        ("Enterprise Architecture.docx", create_docx, FileType.DOCX),
        ("Product Strategy.pptx", create_pptx, FileType.PPTX),
        ("employees.xlsx", create_xlsx, FileType.XLSX),
    ],
)
def test_office_ingestion_is_stable(tmp_path, filename, creator, file_type):
    path = tmp_path / filename
    creator(path)
    service = IngestionService(max_chars=500)
    first = service.ingest(path)
    second = service.ingest(path)
    assert first.document.file_type == file_type
    assert first.document.id == second.document.id
    assert [chunk.id for chunk in first.chunks] == [chunk.id for chunk in second.chunks]
    assert all(chunk.metadata.source_locator for chunk in first.chunks)


def test_docx_extracts_heading_paragraphs_and_semantic_table(tmp_path):
    path = tmp_path / "Enterprise Architecture.docx"
    create_docx(path)
    result = IngestionService(max_chars=500).ingest(path)
    text = "\n".join(chunk.text for chunk in result.chunks)
    assert "Enterprise Architecture" in text
    assert "Apache Spark" in text
    assert "System: Orion" in text
    assert "Owner: Data Platform" in text
    assert any(chunk.metadata.heading == "Enterprise Architecture" for chunk in result.chunks)
    assert any(chunk.metadata.source_locator.type == "table" for chunk in result.chunks)


def test_pptx_extracts_slides_and_table_with_slide_locator(tmp_path):
    path = tmp_path / "Product Strategy.pptx"
    create_pptx(path)
    result = IngestionService(max_chars=500).ingest(path)
    text = "\n".join(chunk.text for chunk in result.chunks)
    assert "Project Atlas" in text
    assert "retrieval-augmented generation" in text
    assert "Owner: Knowledge Systems" in text
    assert {chunk.metadata.source_locator.label for chunk in result.chunks} == {"Slide 1", "Slide 2"}


def test_xlsx_extracts_sheets_headers_rows_and_ignores_blanks(tmp_path):
    path = tmp_path / "employees.xlsx"
    create_xlsx(path)
    result = IngestionService(max_chars=1000).ingest(path)
    text = "\n".join(chunk.text for chunk in result.chunks)
    assert "Name: Jordan Lee" in text
    assert "Skills: PyTorch, RAG" in text
    assert "Sheet: Systems" in text
    assert "\n\n\n" not in text
    assert any("Employees · rows 2–4" == chunk.metadata.source_locator.label for chunk in result.chunks)


def test_xlsx_formula_is_preserved_and_never_calculated(tmp_path):
    path = tmp_path / "formulas.xlsx"
    create_xlsx(path, formula=True)
    _, units = XLSXLoader().load(path)
    text = "\n".join(unit.text for unit in units)
    assert "[Formula: =1+1]" in text
    assert "Dependency: 2" not in text


def test_xlsx_without_dimension_metadata_extracts_rows_and_provenance(tmp_path):
    path = tmp_path / "dimensionless.xlsx"
    create_dimensionless_xlsx(path)

    _, units = XLSXLoader().load(path)

    assert len(units) == 1
    assert "Name: Aster" in units[0].text
    assert units[0].section == "Employees"
    assert units[0].source_locator.label == "Employees · rows 2–2"


def test_xlsx_bounds_are_explicit(tmp_path, monkeypatch):
    path = tmp_path / "employees.xlsx"
    create_xlsx(path)
    monkeypatch.setattr("app.ingestion.loaders.xlsx_loader.settings.xlsx_max_rows_per_sheet", 2)
    with pytest.raises(ValueError, match="row limit"):
        XLSXLoader().load(path)


@pytest.mark.parametrize("suffix", ["docx", "pptx", "xlsx"])
@pytest.mark.parametrize("content", [b"", b"not a zip archive"])
def test_office_loaders_reject_empty_or_malformed_input(tmp_path, suffix, content):
    path = tmp_path / f"broken.{suffix}"
    path.write_bytes(content)
    with pytest.raises(ValueError, match="empty|Malformed"):
        IngestionService().ingest(path)
